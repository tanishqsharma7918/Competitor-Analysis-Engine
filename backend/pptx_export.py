from pptx import Presentation
from pptx.util import Inches, Pt
import io
import streamlit as st


@st.cache_data(show_spinner=False)
def generate_powerpoint_report(
    product_name: str,
    matrix_df,
    analysis,
    features
) -> bytes:
    """
    Full PPT report with:
    - Title
    - Competitor Summary
    - Feature Summary
    - Differentiators
    - Recommendations
    - Gaps
    """

    prs = Presentation()

    # ---------------------------------------------------
    # Title Slide
    # ---------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"Competitor Analysis — {product_name}"
    slide.placeholders[1].text = "AI-Generated Market Intelligence Report"

    # ---------------------------------------------------
    # Feature Summary Slide
    # ---------------------------------------------------
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.placeholders[1]

    title.text = "Feature Summary"
    body.text = f"{len(features)} features extracted and categorized."

    # ---------------------------------------------------
    # Differentiators Slide
    # ---------------------------------------------------
    diffs = analysis.get("differentiators", [])
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Differentiators"

    tf = slide.placeholders[1].text_frame
    for d in diffs:
        p = tf.add_paragraph()
        p.text = f"{d.get('title')}: {d.get('description')}"
        p.level = 1

    # ---------------------------------------------------
    # Recommendations Slide
    # ---------------------------------------------------
    recs = analysis.get("recommendations", [])
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Strategic Recommendations"

    tf = slide.placeholders[1].text_frame
    for r in recs:
        p = tf.add_paragraph()
        p.text = f"{r.get('title')}: {r.get('description')}"
        p.level = 1

    # ---------------------------------------------------
    # Gaps Slide
    # ---------------------------------------------------
    gaps = analysis.get("gaps", [])
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Product Gaps"

    tf = slide.placeholders[1].text_frame
    for g in gaps:
        p = tf.add_paragraph()
        p.text = f"{g.get('title')}: {g.get('description')}"
        p.level = 1

    # ---------------------------------------------------
    # Export
    # ---------------------------------------------------
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
