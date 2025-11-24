import streamlit as st
import pandas as pd
from typing import Dict, Any, List
import io
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt
from fpdf import FPDF


# ---------------------------------------------------------
# STABLE CACHE KEY BUILDER
# ---------------------------------------------------------
def _report_cache_key(product_name: str) -> str:
    return product_name.strip().lower()


# ---------------------------------------------------------
# EXCEL REPORT (CACHED)
# ---------------------------------------------------------
@st.cache_data(show_spinner=False)
def generate_excel_report(product_name: str, df: pd.DataFrame) -> bytes:
    """
    Creates Excel report and caches the result.
    """

    wb = Workbook()
    ws = wb.active
    ws.title = "Comparison Matrix"

    # Write header
    ws.append(["Product"] + list(df.columns))

    # Write rows
    for product, row in df.iterrows():
        ws.append([product] + list(row.values))

    # Save to bytes
    buffer = io.BytesIO()
    wb.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------
# POWERPOINT REPORT (CACHED)
# ---------------------------------------------------------
@st.cache_data(show_spinner=False)
def generate_powerpoint_report(product_name: str, df: pd.DataFrame) -> bytes:
    """
    Creates PPT report and caches output.
    """

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = f"Competitor Analysis — {product_name}"
    subtitle.text = "Auto-generated Report"

    # Add table slide
    slide_layout = prs.slide_layouts[5]
    table_slide = prs.slides.add_slide(slide_layout)

    rows = df.shape[0] + 1
    cols = df.shape[1] + 1

    table = table_slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1), Inches(9), Inches(5)).table

    # Header
    table.cell(0, 0).text = "Product"
    for j, col in enumerate(df.columns, start=1):
        table.cell(0, j).text = col

    # Rows
    for i, (product, row) in enumerate(df.iterrows(), start=1):
        table.cell(i, 0).text = str(product)
        for j, val in enumerate(row.values, start=1):
            table.cell(i, j).text = str(val)

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ---------------------------------------------------------
# PDF REPORT (CACHED)
# ---------------------------------------------------------
@st.cache_data(show_spinner=False)
def generate_pdf_report(product_name: str, df: pd.DataFrame) -> bytes:
    """
    Generates PDF report with caching for speed.
    """

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=14)
    pdf.cell(0, 10, f"Competitor Analysis — {product_name}", ln=True)

    pdf.set_font("Arial", size=10)

    # Table header
    header = ["Product"] + list(df.columns)
    for h in header:
        pdf.cell(40, 10, h, 1)
    pdf.ln()

    # Table rows
    for product, row in df.iterrows():
        pdf.cell(40, 10, str(product), 1)
        for val in row.values:
            pdf.cell(40, 10, str(val), 1)
        pdf.ln()

    # Save to bytes
    buffer = io.BytesIO()
    buffer.write(pdf.output(dest="S").encode("latin-1"))
    return buffer.getvalue()
