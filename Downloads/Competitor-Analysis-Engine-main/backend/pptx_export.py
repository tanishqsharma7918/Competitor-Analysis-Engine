from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List, Dict, Any
import pandas as pd
from io import BytesIO


def generate_powerpoint_report(
    product_name: str,
    competitors: List[Dict[str, Any]],
    matrix_df: pd.DataFrame,
    analysis: Dict[str, Any],
    features: List[Dict[str, Any]]
) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"Competitor Analysis"
    subtitle.text = f"{product_name}\n\nAI-Powered Competitive Intelligence Report"
    
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Executive Summary"
    
    tf = body_shape.text_frame
    tf.text = f"Product Analyzed: {product_name}"
    
    p = tf.add_paragraph()
    p.text = f"Competitors Identified: {len(competitors)}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Features Analyzed: {len(features)}"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = f"Key Differentiators: {len(analysis.get('differentiators', []))}"
    p.level = 1
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Competitors Overview"
    
    tf = body_shape.text_frame
    tf.clear()
    
    for i, comp in enumerate(competitors[:8], 1):
        p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
        p.text = f"{comp.get('company_name', 'Unknown')}: {comp.get('product_name', 'Unknown')}"
        p.level = 0
        
        p_detail = tf.add_paragraph()
        p_detail.text = f"{comp.get('market_position', 'Unknown')} - {comp.get('description', '')[:80]}..."
        p_detail.level = 1
        p_detail.font.size = Pt(11)
    
    differentiators = analysis.get('differentiators', [])
    if differentiators:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Key Differentiators"
        
        tf = body_shape.text_frame
        tf.clear()
        
        for i, diff in enumerate(differentiators[:6], 1):
            p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
            p.text = diff.get('title', '')
            p.level = 0
            p.font.bold = True
            p.font.size = Pt(14)
            
            p_detail = tf.add_paragraph()
            p_detail.text = diff.get('description', '')[:150] + ("..." if len(diff.get('description', '')) > 150 else "")
            p_detail.level = 1
            p_detail.font.size = Pt(11)
    
    recommendations = analysis.get('recommendations', [])
    if recommendations:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Strategic Recommendations"
        
        tf = body_shape.text_frame
        tf.clear()
        
        for i, rec in enumerate(recommendations[:6], 1):
            p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
            p.text = rec.get('title', '')
            p.level = 0
            p.font.bold = True
            p.font.size = Pt(14)
            
            p_detail = tf.add_paragraph()
            p_detail.text = rec.get('description', '')[:150] + ("..." if len(rec.get('description', '')) > 150 else "")
            p_detail.level = 1
            p_detail.font.size = Pt(11)
    
    missing = analysis.get('missing_capabilities', [])
    if missing:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Missing Capabilities"
        
        tf = body_shape.text_frame
        tf.clear()
        
        for i, miss in enumerate(missing[:6], 1):
            p = tf.add_paragraph() if i > 1 else tf.paragraphs[0]
            importance = miss.get('importance', 'Unknown')
            p.text = f"{miss.get('capability', '')} ({importance} Priority)"
            p.level = 0
            p.font.bold = True
            p.font.size = Pt(14)
            
            if importance == 'High':
                p.font.color.rgb = RGBColor(255, 68, 68)
            elif importance == 'Medium':
                p.font.color.rgb = RGBColor(255, 153, 68)
            
            p_detail = tf.add_paragraph()
            p_detail.text = miss.get('rationale', '')[:150] + ("..." if len(miss.get('rationale', '')) > 150 else "")
            p_detail.level = 1
            p_detail.font.size = Pt(11)
    
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.75)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.text = "Feature Comparison Matrix"
    p = text_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    
    note_top = Inches(1.5)
    note_box = slide.shapes.add_textbox(left, note_top, width, Inches(0.5))
    note_frame = note_box.text_frame
    note_frame.text = "Detailed feature matrix available in Excel/PDF reports due to formatting constraints"
    note_frame.paragraphs[0].font.size = Pt(14)
    note_frame.paragraphs[0].font.italic = True
    
    table_top = Inches(2.5)
    rows = min(10, len(matrix_df)) + 1
    cols = min(5, len(matrix_df.columns))
    
    table = slide.shapes.add_table(rows, cols, left, table_top, Inches(9), Inches(4)).table
    
    for col_idx in range(cols):
        cell = table.cell(0, col_idx)
        cell.text = str(matrix_df.columns[col_idx])
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(54, 96, 146)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    for row_idx in range(min(9, len(matrix_df))):
        for col_idx in range(cols):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(matrix_df.iloc[row_idx, col_idx])
            cell.text_frame.paragraphs[0].font.size = Pt(9)
    
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
